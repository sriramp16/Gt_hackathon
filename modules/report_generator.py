import logging
from typing import Dict, List, Any, Optional
from datetime import datetime
from pathlib import Path
import io
import tempfile

try:
    import matplotlib.pyplot as plt
    HAS_MATPLOTLIB = True
except Exception:
    HAS_MATPLOTLIB = False

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image
    from reportlab.platypus import PageBreak, KeepTogether
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    HAS_PYTHON_PPTX = True
except ImportError:
    HAS_PYTHON_PPTX = False


class PDFReportGenerator:
    
    def __init__(self, title: str = "Executive Insight Report"):
        if not HAS_REPORTLAB:
            raise ImportError("reportlab package not installed")
        
        self.title = title
        self.elements = []
        self.styles = getSampleStyleSheet()
        self._setup_custom_styles()
    
    def _setup_custom_styles(self):
        
        self.styles.add(ParagraphStyle(
            name='CustomTitle',
            parent=self.styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1a3a52'),
            spaceAfter=30,
            alignment=1
        ))
        
        self.styles.add(ParagraphStyle(
            name='CustomHeading2',
            parent=self.styles['Heading2'],
            fontSize=14,
            textColor=colors.HexColor('#2d5a7b'),
            spaceAfter=12
        ))
        
        self.styles.add(ParagraphStyle(
            name='Insight',
            parent=self.styles['BodyText'],
            fontSize=12,
            leading=16,
            leftIndent=16,
            spaceAfter=12
        ))
    
    def add_title(self, title: str):
        self.elements.append(Paragraph(title, self.styles['CustomTitle']))
        self.elements.append(Spacer(1, 0.2*inch))
    
    def add_section(self, title: str):
        self.elements.append(Paragraph(title, self.styles['CustomHeading2']))
        self.elements.append(Spacer(1, 0.1*inch))
    
    def add_paragraph(self, text: str, style: str = 'BodyText'):
        self.elements.append(Paragraph(text, self.styles[style]))
        self.elements.append(Spacer(1, 0.1*inch))
    
    def add_bullet_points(self, points: List[str]):
        for point in points:
            self.elements.append(Paragraph(f"• {point}", self.styles['Insight']))
        self.elements.append(Spacer(1, 0.1*inch))
    
    def add_table(self, data: List[List], col_widths: Optional[List] = None):
        if col_widths is None:
            col_widths = [1.5*inch] * len(data[0]) if data else []
        
        table = Table(data, colWidths=col_widths)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2d5a7b')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 11),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        
        self.elements.append(table)
        self.elements.append(Spacer(1, 0.2*inch))
    
    def add_page_break(self):
        self.elements.append(PageBreak())
    
    def generate(self, output_path: str):
        doc = SimpleDocTemplate(
            output_path,
            pagesize=letter,
            rightMargin=0.5*inch,
            leftMargin=0.5*inch,
            topMargin=0.75*inch,
            bottomMargin=0.75*inch
        )
        
        doc.build(self.elements)
        logger.info(f"PDF report generated: {output_path}")


class PowerPointReportGenerator:
    
    def __init__(self, title: str = "Executive Insight Report"):
        if not HAS_PYTHON_PPTX:
            raise ImportError("python-pptx package not installed")
        
        self.presentation = Presentation()
        self.presentation.slide_width = Inches(10)
        self.presentation.slide_height = Inches(7.5)
        self.title = title
        self._setup_master()
    
    def _setup_master(self):
        pass
    
    def add_title_slide(self, title: str, subtitle: str = ""):
        slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[6]  # Blank layout
        )
        
        # Add background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(26, 58, 82)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(9), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4), Inches(9), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_frame.paragraphs[0].font.size = Pt(24)
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(220, 220, 220)
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add date
        date_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.5), Inches(9), Inches(0.5)
        )
        date_frame = date_box.text_frame
        date_frame.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
        date_frame.paragraphs[0].font.size = Pt(12)
        date_frame.paragraphs[0].font.color.rgb = RGBColor(180, 180, 180)
        date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def add_content_slide(self, title: str, content: Dict[str, Any]):
        slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[6]  # Blank layout
        )
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(45, 90, 123)
        
        # Add separator line
        slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(9), Inches(0))
        
        # Add content
        if "text" in content:
            text_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.7)
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.text = content["text"]
            text_frame.paragraphs[0].font.size = Pt(14)
        
        if "bullets" in content:
            bullet_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.7)
            )
            bullet_frame = bullet_box.text_frame
            bullet_frame.word_wrap = True
            
            for i, bullet in enumerate(content["bullets"]):
                if i == 0:
                    p = bullet_frame.paragraphs[0]
                else:
                    p = bullet_frame.add_paragraph()
                
                p.text = bullet
                p.font.size = Pt(14)
                p.level = 0
                p.space_before = Pt(6)
    
    def add_table_slide(self, title: str, table_data: List[List[str]]):
        """Add slide with table"""
        
        slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[6]
        )
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(45, 90, 123)
        
        # Add table
        rows, cols = len(table_data), len(table_data[0]) if table_data else 0
        left = Inches(0.5)
        top = Inches(1.3)
        width = Inches(9)
        height = Inches(5.5)
        
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        for i, row in enumerate(table_data):
            for j, cell_value in enumerate(row):
                cell = table_shape.cell(i, j)
                cell.text = str(cell_value)
                
                # Format header row
                if i == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(45, 90, 123)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(255, 255, 255)
    
    def generate(self, output_path: str):
        self.presentation.save(output_path)
        logger.info(f"PowerPoint report generated: {output_path}")


class ReportBuilder:
    
    def __init__(self, analysis_results: Dict[str, Any], insights: Dict[str, str]):
        self.analysis = analysis_results
        self.insights = insights
    
    def generate_pdf_report(self, output_path: str, title: str = "Executive Insight Report"):
        pdf = PDFReportGenerator(title)
        
        # Title page
        pdf.add_title(title)
        pdf.add_paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y')}")
        pdf.add_page_break()
        
        # Executive Summary (use bullets if possible)
        pdf.add_section("Executive Summary")
        if "kpi_summary" in self.insights:
            kpi_text = self.insights["kpi_summary"] or ""
            bullets = [s.strip() for s in kpi_text.split('\n') if s.strip()]
            if len(bullets) > 1:
                pdf.add_bullet_points(bullets)
            else:
                pdf.add_paragraph(kpi_text, style='Insight')
        pdf.add_page_break()

        # Key Findings
        pdf.add_section("Key Findings")
        if "performance_analysis" in self.insights:
            perf_text = self.insights["performance_analysis"] or ""
            perf_bullets = [s.strip() for s in perf_text.split('\n') if s.strip()]
            if len(perf_bullets) > 1:
                pdf.add_bullet_points(perf_bullets)
            else:
                pdf.add_paragraph(perf_text, style='Insight')
        pdf.add_page_break()

        # Recommendations
        pdf.add_section("Recommendations")
        if "recommendations" in self.insights:
            rec_text = self.insights["recommendations"] or ""
            rec_bullets = [s.strip() for s in rec_text.split('\n') if s.strip()]
            if len(rec_bullets) > 0:
                pdf.add_bullet_points(rec_bullets[:8])
            else:
                pdf.add_paragraph(rec_text, style='Insight')
        pdf.add_page_break()

        # Data Quality
        if "quality_report" in self.analysis:
            pdf.add_section("Data Quality Report")
            pdf.add_bullet_points([
                f"Data shape: {self.analysis['quality_report']}",
                "All data processed and validated"
            ])

        # KPI Table (if available)
        if "kpis" in self.analysis and isinstance(self.analysis["kpis"], dict):
            from modules.utils import TextFormatter

            pdf.add_section("Key Metrics")
            kpi_items = self.analysis["kpis"]
            table_data = [["Metric", "Value"]]
            for k, v in list(kpi_items.items())[:12]:
                # Format numbers and percentages nicely
                if isinstance(v, (int, float)):
                    formatted = TextFormatter.format_number(v, decimals=2)
                elif isinstance(v, dict):
                    # For nested dicts, show brief summary (e.g., counts)
                    try:
                        formatted = ", ".join([f"{kk}: {vv}" for kk, vv in list(v.items())[:3]])
                    except Exception:
                        formatted = str(v)
                else:
                    # Try to coerce numeric-like strings
                    try:
                        fv = float(str(v))
                        formatted = TextFormatter.format_number(fv, decimals=2)
                    except Exception:
                        formatted = str(v)

                table_data.append([str(k), formatted])

            pdf.add_table(table_data, col_widths=[3*inch, 3*inch])
            pdf.add_page_break()

        # KPI Chart (improved visuals, annotations, single-item handling)
        if HAS_MATPLOTLIB and "kpis" in self.analysis and isinstance(self.analysis["kpis"], dict):
            try:
                # build numeric dict (coerce strings if possible)
                numeric_items = {}
                for k, v in self.analysis["kpis"].items():
                    try:
                        # skip nested dicts
                        if isinstance(v, dict):
                            continue
                        numeric_items[k] = float(v)
                    except Exception:
                        # ignore non-numeric
                        continue

                # keep only numeric values
                numeric_items = {k: v for k, v in numeric_items.items() if v is not None}

                if len(numeric_items) > 0:
                    # sort by value desc and take top 6 (or fewer)
                    items = sorted(numeric_items.items(), key=lambda x: x[1], reverse=True)[:6]
                    labels = [str(x[0]) for x in items]
                    values = [x[1] for x in items]

                    # If only one KPI, use horizontal bar for better appearance
                    figsize = (6, 3) if len(items) > 1 else (8, 2.5)
                    fig, ax = plt.subplots(figsize=figsize)

                    if len(items) == 1:
                        ax.barh(labels, values, color="#2d5a7b")
                        ax.invert_yaxis()
                        for i, v in enumerate(values):
                            ax.text(v + max(values) * 0.01, i, f"{v:,.0f}", va='center', fontsize=10)
                        ax.set_xlabel("Count")
                    else:
                        bars = ax.bar(labels, values, color="#2d5a7b")
                        for bar in bars:
                            h = bar.get_height()
                            ax.annotate(f"{h:,.0f}", xy=(bar.get_x() + bar.get_width() / 2, h),
                                        xytext=(0, 6), textcoords="offset points", ha='center', fontsize=10)
                        ax.set_ylabel("Count")
                        ax.set_xticklabels(labels, rotation=30, ha='right')

                    ax.set_title("Top KPIs")
                    ax.grid(axis='y', linestyle='--', alpha=0.4)
                    plt.tight_layout()

                    buf = io.BytesIO()
                    fig.savefig(buf, format='png', dpi=200, bbox_inches='tight')
                    plt.close(fig)
                    buf.seek(0)

                    img_temp = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                    img_temp.write(buf.read())
                    img_temp.flush()
                    img_temp.close()

                    pdf.elements.append(Image(img_temp.name, width=6*inch, height=3*inch))
                    pdf.add_page_break()
            except Exception as e:
                logger.warning(f"KPI chart generation failed: {e}")

        pdf.generate(output_path)
        logger.info(f"PDF report saved to {output_path}")
    
    def generate_pptx_report(self, output_path: str, title: str = "Executive Insight Report"):
        pptx = PowerPointReportGenerator(title)
        
        # Title slide
        pptx.add_title_slide(title, "Automated Insight Analysis")
        
        # Executive Summary
        pptx.add_content_slide("Executive Summary", {
            "text": self.insights.get("kpi_summary", "Analysis complete")
        })
        
        # Key Findings
        if "performance_analysis" in self.insights:
            pptx.add_content_slide("Key Findings", {
                "text": self.insights["performance_analysis"]
            })
        
        # Recommendations
        if "recommendations" in self.insights:
            recommendations_text = self.insights["recommendations"]
            bullets = [line.strip() for line in recommendations_text.split('\n') 
                      if line.strip() and not line.startswith('#')]
            
            pptx.add_content_slide("Recommendations", {
                "bullets": bullets[:5]  # Limit to 5 bullets per slide
            })
        
        # Data Insights
        pptx.add_content_slide("About This Report", {
            "bullets": [
                "Automated analysis using advanced data processing",
                "AI-powered insights generation",
                f"Generated: {datetime.now().strftime('%B %d, %Y')}",
                "For questions, contact your data analytics team"
            ]
        })
        
        pptx.generate(output_path)
        logger.info(f"PowerPoint report saved to {output_path}")
